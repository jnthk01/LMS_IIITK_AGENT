import os
import pytesseract
from PIL import Image
from io import BytesIO
import streamlit as st
from pptx import Presentation
from langchain_chroma import Chroma
from langchain.prompts import ChatPromptTemplate 
from rag_file_retirever import rag_file_retriever
from langchain.chains import create_retrieval_chain
from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI

def rag_streamlit():

    def extract_text_from_pptx(file_path):
        prs = Presentation(file_path)
        text_data = []
        
        for slide_number, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text.strip())
            text_data.append(f"Slide {slide_number + 1}:\n" + "\n".join(slide_text))

        return "\n\n".join(text_data)


    def extract_text_from_images_sub(slide):
        image_text = []
        for shape in slide.shapes:
            if shape.shape_type == 13:  
                image_stream = shape.image.blob  
                image = Image.open(BytesIO(image_stream))
                extracted_text = pytesseract.image_to_string(image)
                if extracted_text.strip():
                    image_text.append(extracted_text.strip())
        return "\n".join(image_text)

    def extract_text_from_images(file_path):
        text_data = []

        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            img_text = extract_text_from_images_sub(slide)
            if img_text:
                text_data.append(f"Slide {i + 1}: with image text: \n" + "\n".join(img_text))
        
        return text_data

    def combine_ppt_text_and_images_by_slide(ppt_text, ppt_img):
        combined_list = []

        ppt_text_dict = {}
        ppt_img_dict = {}

        for entry in ppt_text:
            slide_number = int(entry.split(":")[0].split()[-1])
            ppt_text_dict[slide_number] = entry.strip()

        for entry in ppt_img:
            slide_number = int(entry.split(":")[0].split()[-1])
            ppt_img_dict[slide_number] = entry.strip()
        all_slide_numbers = sorted(set(ppt_text_dict.keys()).union(set(ppt_img_dict.keys())))

        for slide_number in all_slide_numbers:
            combined_content = ppt_text_dict.get(slide_number, "") + "\n\n" + ppt_img_dict.get(slide_number, "")
            if combined_content.strip():
                combined_list.append(f"Slide {slide_number}:\n{combined_content.strip()}")

        return combined_list

    def extract_text_from_pdf(file_path):
        loader = PyPDFLoader(file_path)
        return loader.load()

    def chunk_text(text):
        splitter = RecursiveCharacterTextSplitter(chunk_size=1200, chunk_overlap=300)
        return splitter.split_documents(text)

    GEMINI_API_KEY = os.getenv('GEMINI_API_KEY')

    def generate_vector_db(file_name, chunks_data):
        vector_db = Chroma(
            collection_name=file_name,
            embedding_function=GoogleGenerativeAIEmbeddings(google_api_key=GEMINI_API_KEY, model="models/text-embedding-004"),
            persist_directory=file_name+"db"
        )
        vector_db.add_documents(chunks_data)
        return vector_db

    def load_or_create_vector_db(file_name, chunks_data):
        vector_db_path = file_name+"db"
        
        if os.path.exists(vector_db_path):
            vector_db = Chroma(
                collection_name=file_name,
                embedding_function=GoogleGenerativeAIEmbeddings(google_api_key=GEMINI_API_KEY, model="models/text-embedding-004"),
                persist_directory=file_name+"db"
            )
        else:
            vector_db = generate_vector_db(file_name, chunks_data)
        
        return vector_db

    st.title("Document Query with LangChain")

    user_type = st.selectbox("Select File Type", ["pdf", "pptx"])
    user_prompt = st.text_input("Enter the URL or File Path")
    file_name_input = st.text_input("Enter File Name")

    if user_prompt and file_name_input:  
        if user_type.lower() == "pdf":
            file_name = rag_file_retriever(user_type=user_type, user_prompt=user_prompt, file_name=file_name_input)
            print("Extracting PDF DATA")
            pdf_text = extract_text_from_pdf(file_name)
            chunks_data = chunk_text(pdf_text)

        elif user_type.lower() == "pptx":
            file_name = rag_file_retriever(user_type=user_type, user_prompt=user_prompt, file_name=file_name_input)
            ppt_text = extract_text_from_pptx(file_name)
            ppt_img = extract_text_from_images(file_name)
            final_ppt_data = combine_ppt_text_and_images_by_slide(ppt_text=ppt_text, ppt_img=ppt_img)
            final_ppt_data = "\n".join(final_ppt_data)
            chunks_data = chunk_text(final_ppt_data)

        print("Vector DB creation")
        vector_db = load_or_create_vector_db(file_name, chunks_data)

        retriever = vector_db.as_retriever(search_type="similarity", search_kwargs={"k": 10})

        llm = ChatGoogleGenerativeAI(
            api_key=GEMINI_API_KEY,
            model="gemini-1.5-flash",
            temperature=0.3,
            max_tokens=None,
            timeout=None
        )

        prompt_template = ChatPromptTemplate.from_template(""" 
        Answer the following question based only on the provided context:
        <context>
        {context}
        </context>
        Question: {input}
        """)

        combine_docs_chain = create_stuff_documents_chain(llm, prompt_template)
        retrieval_chain = create_retrieval_chain(retriever, combine_docs_chain)

        user_query = st.text_input("Ask a question about the document")

        if user_query:
            response = retrieval_chain.invoke({"input": user_query})
            st.write(response['answer'])